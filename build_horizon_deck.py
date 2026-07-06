"""
build_horizon_deck.py
=====================
Genera un deck de 3 slides para Horizon Program con branding Alkalma:
  1) Cover — "Horizon Program" + fecha
  2) Dos QR de surveys (Horizon Workshops Registration + Customer Satisfaction)
  3) Patient Feedback Journey (imagen full-width)

Importa los helpers de branding del skill oficial Alkalma para reutilizar
paleta, tipografía, gradiente y logo exactos del brandbook.
"""

from __future__ import annotations

import sys
from pathlib import Path

SKILL_SCRIPTS = Path.home() / ".claude" / "skills" / "branding_web_alkalma" / "scripts"
sys.path.insert(0, str(SKILL_SCRIPTS))

from build_alkalma_pptx import (  # type: ignore
    INDIGO, STEEL_BLUE, TURQUOISE, BLUEGRAY, AZURE_WHITE, WHITE,
    SLIDE_W, SLIDE_H,
    add_gradient_background, add_solid_background, add_radial_overlay,
    add_text, add_logo, add_header, add_footer, add_chip,
    _lighten_for_bg,
)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

# Rutas de las imágenes
QR_HORIZON = Path("/Users/danielsandoval/Downloads/QRCode for Horizon Workshops Registration Form.png")
QR_SATISFACTION = Path("/Users/danielsandoval/Downloads/QRCode for Customer Satisfaction Survey.png")
PATIENT_JOURNEY = Path("/Users/danielsandoval/Downloads/alkalma_tmp/patient survay.png")

OUTPUT = Path("/Users/danielsandoval/Documents/Proyectos/test_claude/Horizon_Program_Deck.pptx")


def build_slide_1_cover(slide):
    """Slide 1: Cover con 'Horizon Program' y la fecha."""
    add_gradient_background(slide, INDIGO, STEEL_BLUE, 135)
    add_radial_overlay(slide, x_in=9.5, y_in=-2.5, size_in=8, color=TURQUOISE, alpha_pct=22)
    add_radial_overlay(slide, x_in=-2.0, y_in=4.0, size_in=6, color=TURQUOISE, alpha_pct=12)

    # Eyebrow chip
    add_chip(
        slide, "UNITED IN WELL-BEING",
        x_in=0.9, y_in=1.7,
        color_bg=RGBColor(0x33, 0x66, 0x80),
        color_text=WHITE,
        w_in=2.8, h_in=0.45,
    )

    # Titulo principal
    add_text(
        slide, "Horizon Program",
        x_in=0.9, y_in=2.5, w_in=11.5, h_in=2.0,
        font_size=84, weight=300, color=WHITE,
        emphasis_text="Horizon",
        emphasis_weight=700,
        line_spacing=1.05,
    )

    # Subtítulo
    add_text(
        slide, "Workshops & patient experience initiative",
        x_in=0.9, y_in=4.7, w_in=10.5, h_in=0.7,
        font_size=20, weight=400,
        color=RGBColor(0xE5, 0xEE, 0xF2),
        line_spacing=1.4,
    )

    # Línea separadora turquoise
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.9), Inches(5.7), Inches(0.8), Inches(0.05),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = TURQUOISE
    line.line.fill.background()

    # Fecha
    add_text(
        slide, "9 de Junio, 2026",
        x_in=0.9, y_in=5.9, w_in=10.5, h_in=0.5,
        font_size=22, weight=500, color=WHITE,
        letter_spacing=1.0,
    )

    # Logo + tagline footer
    add_logo(slide, x_in=0.9, y_in=6.6, height_in=0.55, white=True)


def build_slide_2_qr(slide):
    """Slide 2: Dos QR codes lado a lado — Registration + Satisfaction Survey."""
    add_solid_background(slide, AZURE_WHITE)
    add_header(slide, page_num="02 / 03")

    # Título de la sección
    add_text(
        slide, "Scan to participate",
        x_in=0.7, y_in=1.4, w_in=12, h_in=0.9,
        font_size=38, weight=500, color=INDIGO,
        emphasis_text="participate",
        emphasis_weight=700, line_spacing=1.1,
    )

    # Subtítulo
    add_text(
        slide, "Two ways to engage with the Horizon Program — register for our workshops or share your experience.",
        x_in=0.7, y_in=2.35, w_in=11.5, h_in=0.7,
        font_size=15, weight=400, color=BLUEGRAY,
        line_spacing=1.5,
    )

    # Dos cards con QR codes
    card_w = 5.4
    card_h = 3.8
    card_y = 3.25
    gap = 0.5
    total_w = card_w * 2 + gap
    start_x = (13.333 - total_w) / 2

    qr_data = [
        {
            "qr_path": QR_HORIZON,
            "title": "Workshops Registration",
            "subtitle": "Sign up for the Horizon Program workshops",
            "icon": "→",
        },
        {
            "qr_path": QR_SATISFACTION,
            "title": "Satisfaction Survey",
            "subtitle": "Share your feedback about your experience",
            "icon": "★",
        },
    ]

    for i, qr in enumerate(qr_data):
        x = start_x + i * (card_w + gap)

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(card_y), Inches(card_w), Inches(card_h),
        )
        card.adjustments[0] = 0.05
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xDC, 0xE5, 0xEA)
        card.line.width = Pt(0.75)
        card.text_frame.text = ""

        # QR code image — centrado en la card, parte superior
        qr_size = 2.3
        qr_x = x + (card_w - qr_size) / 2
        qr_y = card_y + 0.35
        if qr["qr_path"].exists():
            slide.shapes.add_picture(
                str(qr["qr_path"]),
                Inches(qr_x), Inches(qr_y),
                width=Inches(qr_size), height=Inches(qr_size),
            )

        # Titulo de la card
        add_text(
            slide, qr["title"],
            x_in=x + 0.3, y_in=card_y + 2.85, w_in=card_w - 0.6, h_in=0.5,
            font_size=20, weight=600, color=INDIGO,
            align="center",
        )

        # Subtitle
        add_text(
            slide, qr["subtitle"],
            x_in=x + 0.4, y_in=card_y + 3.3, w_in=card_w - 0.8, h_in=0.5,
            font_size=12, weight=400, color=BLUEGRAY,
            align="center", line_spacing=1.4,
        )

    add_footer(slide, left_text="alkalma · United in Well-Being", right_text="HORIZON PROGRAM")


def build_slide_3_journey(slide):
    """Slide 3: Patient Feedback Journey — imagen full-width."""
    add_solid_background(slide, AZURE_WHITE)
    add_header(slide, page_num="03 / 03")

    # Título
    add_text(
        slide, "Patient Feedback Journey",
        x_in=0.7, y_in=1.4, w_in=12, h_in=0.9,
        font_size=38, weight=500, color=INDIGO,
        emphasis_text="Journey",
        emphasis_weight=700, line_spacing=1.1,
    )

    # Subtítulo
    add_text(
        slide, "We're listening. Multiple touchpoints designed to capture your voice and improve your care experience.",
        x_in=0.7, y_in=2.35, w_in=11.5, h_in=0.7,
        font_size=15, weight=400, color=BLUEGRAY,
        line_spacing=1.5,
    )

    # Imagen del journey - full width respetando aspect ratio (2216 x 760 ≈ 2.92:1)
    if PATIENT_JOURNEY.exists():
        img_w = 12.0
        img_h = img_w / (2216 / 760)  # ≈ 4.11
        img_x = (13.333 - img_w) / 2
        img_y = 3.15
        slide.shapes.add_picture(
            str(PATIENT_JOURNEY),
            Inches(img_x), Inches(img_y),
            width=Inches(img_w), height=Inches(img_h),
        )

    add_footer(slide, left_text="alkalma · United in Well-Being", right_text="HORIZON PROGRAM")


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]

    s1 = prs.slides.add_slide(blank_layout)
    build_slide_1_cover(s1)

    s2 = prs.slides.add_slide(blank_layout)
    build_slide_2_qr(s2)

    s3 = prs.slides.add_slide(blank_layout)
    build_slide_3_journey(s3)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"OK Generado: {OUTPUT}")


if __name__ == "__main__":
    main()
